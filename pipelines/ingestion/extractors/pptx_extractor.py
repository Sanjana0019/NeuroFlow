from pathlib import Path

from pptx import Presentation

from pipelines.ingestion.models import ExtractedPage


class PPTXExtractor:
    """Extract slide text content and speaker notes from PowerPoint presentations."""

    def extract(self, file_path: str | Path) -> list[ExtractedPage]:
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {path}")

        if not path.is_file():
            raise ValueError(f"PPTX path is not a file: {path}")

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ValueError(
                f"Failed to read PPTX file '{path}': {exc}"
            ) from exc

        pages: list[ExtractedPage] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []

            # Extract text from slide shapes
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                text = shape.text.strip()
                if text:
                    text_parts.append(text)

            # Extract speaker notes if present
            notes_text = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()

            content_lines = list(text_parts)
            if notes_text:
                content_lines.append(f"Speaker Notes:\n{notes_text}")

            content = "\n\n".join(content_lines).strip()

            if not content:
                continue

            metadata = {
                "source": "pptx",
                "filename": path.name,
                "slide_number": slide_number,
                "has_notes": bool(notes_text),
            }

            pages.append(
                ExtractedPage(
                    page_number=slide_number,
                    content=content,
                    content_type="text",
                    metadata=metadata,
                )
            )

        return pages