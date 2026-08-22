from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from pipelines.ingestion.extractors.pptx_extractor import PPTXExtractor
from pipelines.ingestion.models import ExtractedPage

FIXTURE_PATH = Path("backend/tests/fixtures/test_presentation.pptx")


def create_fixture():
    presentation = Presentation()

    # Slide 1: without notes
    slide1 = presentation.slides.add_slide(
        presentation.slide_layouts[5]
    )
    textbox1 = slide1.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(2),
    )
    textbox1.text = "NeuroFlow PPTX Test"

    # Slide 2: with speaker notes
    slide2 = presentation.slides.add_slide(
        presentation.slide_layouts[5]
    )
    textbox2 = slide2.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(2),
    )
    textbox2.text = "Task 4 Presentation Extraction"

    # Add notes to slide 2
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "Remember to emphasize vector embeddings and pgvector storage."

    presentation.save(FIXTURE_PATH)


def test_pptx_extractor():
    create_fixture()

    extractor = PPTXExtractor()
    pages = extractor.extract(FIXTURE_PATH)

    assert len(pages) == 2

    first_page = pages[0]
    second_page = pages[1]

    # Slide 1
    assert isinstance(first_page, ExtractedPage)
    assert first_page.page_number == 1
    assert first_page.content_type == "text"
    assert "NeuroFlow PPTX Test" in first_page.content
    assert first_page.metadata["source"] == "pptx"
    assert first_page.metadata["filename"] == "test_presentation.pptx"
    assert first_page.metadata["slide_number"] == 1
    assert first_page.metadata["has_notes"] is False

    # Slide 2 with speaker notes
    assert isinstance(second_page, ExtractedPage)
    assert second_page.page_number == 2
    assert "Task 4 Presentation Extraction" in second_page.content
    assert "Speaker Notes:" in second_page.content
    assert "Remember to emphasize vector embeddings" in second_page.content
    assert second_page.metadata["slide_number"] == 2
    assert second_page.metadata["has_notes"] is True